"""Gera os slides PPTX do seminário. Execute: python gerar_apresentacao.py."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("apresentacao_coverage.pptx")
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(12, 25, 46)
BLUE = RGBColor(35, 116, 204)
CYAN = RGBColor(53, 205, 222)
MINT = RGBColor(56, 211, 159)
WHITE = RGBColor(248, 250, 252)
INK = RGBColor(24, 35, 52)
MUTED = RGBColor(90, 108, 130)
PALE = RGBColor(236, 243, 250)
RED = RGBColor(234, 87, 87)
AMBER = RGBColor(244, 180, 57)


def box(slide, x, y, w, h, color, radius=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def text(slide, value, x, y, w, h, size=20, color=INK, bold=False,
         align=PP_ALIGN.LEFT, font="Aptos", valign=MSO_ANCHOR.MIDDLE):
    # Atalhos usados nos blocos de código: o último argumento pode ser a fonte.
    if isinstance(align, str):
        font, align = align, PP_ALIGN.LEFT
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = value
    run.font.name = font; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    tf.margin_left = tf.margin_right = Inches(0.04)
    return tb


def rich(slide, lines, x, y, w, h, size=17, color=INK):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(.05)
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(11)
        r = p.add_run(); r.text = item
        r.font.name = "Aptos"; r.font.size = Pt(size); r.font.color.rgb = color
    return tb


def base(number, section):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    box(s, 0, 0, 13.333, 7.5, WHITE, False)
    box(s, 0, 0, 13.333, .16, CYAN, False)
    text(s, section.upper(), .55, .3, 9, .35, 10, BLUE, True)
    text(s, f"{number:02d}", 12.1, .27, .6, .35, 11, MUTED, True, PP_ALIGN.RIGHT)
    text(s, "PC III · Tema 8 · Gabriel Martins Nunes & Rafaela Garcia Bernardes", .55, 7.05, 10.8, .25, 9, MUTED)
    return s


def title(s, value, subtitle=None):
    text(s, value, .55, .74, 12.1, .65, 29, NAVY, True)
    if subtitle: text(s, subtitle, .57, 1.4, 11.8, .36, 13, MUTED)


# 1 — cover
s = prs.slides.add_slide(prs.slide_layouts[6]); box(s, 0, 0, 13.333, 7.5, NAVY, False)
box(s, .62, .68, 1.3, .12, CYAN, False); text(s, "PROGRAMAÇÃO DE COMPUTADORES III", .65, .93, 7, .35, 13, CYAN, True)
text(s, "Métricas e\nCobertura de Código", .62, 1.55, 7.6, 1.65, 34, WHITE, True)
text(s, "Como medir a abrangência dos testes —\ne por que 100% não significa software sem bugs.", .65, 3.45, 6.6, .8, 17, RGBColor(198, 215, 235))
box(s, 8.35, 1.2, 3.65, 4.55, RGBColor(20, 43, 74))
text(s, "coverage\nreport", 8.75, 1.72, 2.85, .8, 25, WHITE, True, PP_ALIGN.CENTER, "Consolas")
text(s, "Stmts   Miss  Branch  BrPart  Cover\n    8      0      6       0   100%", 8.7, 3.05, 2.9, .8, 13, CYAN, False, PP_ALIGN.CENTER, "Consolas")
box(s, 8.92, 4.32, 2.45, .55, MINT); text(s, "TESTES EXECUTADOS", 9.02, 4.4, 2.25, .22, 10, NAVY, True, PP_ALIGN.CENTER)
text(s, "Gabriel Martins Nunes\nRafaela Garcia Bernardes\nG08PCOM3.01 · T01 · 2026.2", .65, 6.15, 6.2, .65, 14, WHITE)

# 2
s=base(2,"Pergunta central"); title(s,"Teste verde é suficiente?","O ponto de partida da nossa investigação.")
box(s,.7,2.05,5.7,3.9,PALE); text(s,"✓  1 teste passou",1.15,2.42,4.6,.48,27,MINT,True)
text(s,"Podemos concluir que as regras\nda aplicação estão bem testadas?",1.15,3.35,4.8,.8,22,NAVY,True)
box(s,7.05,2.05,5.55,3.9,NAVY); text(s,"Não necessariamente.",7.55,2.45,4.55,.48,27,WHITE,True)
text(s,"O teste pode percorrer apenas\num caminho e ignorar decisões,\nerros e valores-limite.",7.55,3.35,4.3,1.2,19,RGBColor(210,225,241))
text(s,"Coverage torna essas lacunas visíveis.",7.55,5.08,4.3,.3,13,CYAN,True)

# 3
s=base(3,"Conceito"); title(s,"O que é Code Coverage?","Métrica de execução, não selo de qualidade.")
for x,head,body,c in [(0.7,"MEDIR","Quais partes do código foram executadas pelos testes.",BLUE),(4.55,"LOCALIZAR","Onde existem linhas ou caminhos ainda não exercitados.",CYAN),(8.4,"ORIENTAR","Quais cenários devem ser acrescentados à suíte.",MINT)]:
 box(s,x,2.2,3.3,2.85,PALE); box(s,x+.25,2.48,.55,.55,c); text(s,head,x+.95,2.45,2.05,.3,15,NAVY,True); text(s,body,x+.34,3.35,2.65,1.05,16,INK)
text(s,"Coverage responde “quanto foi executado?”. Bons testes respondem “o comportamento correto foi verificado?”.",.85,5.75,11.5,.5,17,NAVY,True,PP_ALIGN.CENTER)

# 4
s=base(4,"Ferramentas"); title(s,"pytest + coverage.py","Ferramentas complementares, com responsabilidades diferentes.")
box(s,.75,2.0,5.75,3.55,NRGB:=RGBColor(237,244,253)); text(s,"pytest",1.18,2.43,2,.45,29,BLUE,True)
rich(s,["• Descobre e executa testes.","• Mostra falhas, erros e testes aprovados.","• Avalia as asserções."],1.2,3.18,4.75,1.65,18)
box(s,6.85,2.0,5.75,3.55,NAVY); text(s,"coverage.py",7.28,2.43,3.3,.45,29,CYAN,True)
rich(s,["• Observa a execução durante os testes.","• Mede linhas e caminhos percorridos.","• Gera relatórios textual e HTML."],7.3,3.18,4.75,1.65,18,RGBColor(220,231,244))
text(s,"coverage run --branch -m pytest",3.55,6.05,6.3,.36,17,NAVY,True,PP_ALIGN.CENTER,"Consolas")

# 5
s=base(5,"Métricas"); title(s,"Statement × Branch coverage","Executar uma linha não é o mesmo que percorrer todas as decisões.")
box(s,.7,2.0,3.4,3.65,PALE); text(s,"if idade >= 18:",1.0,2.42,2.8,.3,18,NAVY,True,"Consolas"); text(s,'    return "adulto"',1.0,2.9,2.8,.3,17,INK,False,"Consolas"); text(s,"else:",1.0,3.38,2.8,.3,18,NAVY,True,"Consolas"); text(s,'    return "menor"',1.0,3.86,2.8,.3,17,INK,False,"Consolas")
box(s,4.55,2.0,3.65,3.65,RGBColor(228,247,243)); text(s,"STATEMENT",4.95,2.47,2.85,.3,16,MINT,True); text(s,"As instruções foram\nexecutadas?",4.95,3.22,2.75,.65,20,NAVY,True); text(s,"idade = 20 executa\napenas o retorno adulto.",4.95,4.38,2.75,.55,15,MUTED)
box(s,8.65,2.0,3.65,3.65,RGBColor(235,241,253)); text(s,"BRANCH",9.05,2.47,2.85,.3,16,BLUE,True); text(s,"Os dois caminhos\nforam percorridos?",9.05,3.22,2.75,.65,20,NAVY,True); text(s,"Precisamos testar\nidade = 20 e idade = 17.",9.05,4.38,2.75,.55,15,MUTED)

# 6
s=base(6,"Demonstração"); title(s,"O experimento em quatro atos","Vamos usar uma regra pequena, com decisões suficientes para gerar evidências.")
steps=[("01","Teste VIP passa","Cobertura inicial"),("02","Ler lacunas","Linhas + branches"),("03","Ampliar a suíte","Limite + exceção"),("04","Questionar o 100%","Teste fraco com bug")]
for i,(n,h,b) in enumerate(steps):
 x=.7+i*3.05; box(s,x,2.45,2.55,2.2,PALE); box(s,x+.25,2.73,.5,.5,[BLUE,CYAN,MINT,AMBER][i]); text(s,n,x+.25,2.82,.5,.16,11,NAVY,True,PP_ALIGN.CENTER); text(s,h,x+.28,3.48,2,.43,16,NAVY,True); text(s,b,x+.28,4.02,2.02,.35,13,MUTED)
 if i<3: text(s,"→",x+2.58,3.3,.35,.3,22,BLUE,True,PP_ALIGN.CENTER)

# 7
s=base(7,"Prática · cenário 1"); title(s,"Um teste verde, três caminhos esquecidos","Gabriel: execute o teste isolado e leia o relatório antes de ampliar a suíte.")
box(s,.7,2.0,5.85,3.6,NAVY); text(s,"def test_cliente_vip():\n    assert calcular_desconto(100, True) == 80",1.05,2.75,5.1,.8,16,WHITE,False,"Consolas")
text(s,"pytest roteiro_live/teste_inicial.py -v",.88,4.82,5.45,.25,12,CYAN,True,"Consolas")
box(s,6.95,2.0,5.65,3.6,PALE); text(s,"O que ainda não foi exercitado?",7.32,2.4,4.9,.35,19,NAVY,True)
rich(s,["• valor negativo → exceção", "• cliente comum abaixo de R$ 500", "• cliente comum a partir de R$ 500"],7.35,3.15,4.6,1.5,17)
text(s,"coverage run --branch -m pytest roteiro_live/teste_inicial.py",.82,6.12,11.75,.3,14,BLUE,True,PP_ALIGN.CENTER,"Consolas")

# 8
s=base(8,"Prática · cenário 2"); title(s,"Novos testes cobrem as regras e os limites","Rafaela: parametrização reduz repetição e deixa os cenários explícitos.")
box(s,.7,1.95,6.3,3.95,PALE); text(s,"@pytest.mark.parametrize(\n    \"valor, vip, esperado\",\n    [(100, True, 80), (100, False, 100),\n     (500, False, 450), (1000, False, 900)]\n)\ndef test_calcular_desconto(valor, vip, esperado):\n    assert calcular_desconto(valor, vip) == esperado",1.02,2.24,5.8,2.65,14,NAVY,False,"Consolas")
box(s,7.42,1.95,5.18,3.95,RGBColor(228,247,243)); text(s,"Além dos valores normais",7.8,2.38,4.3,.3,18,NAVY,True); text(s,"Teste a fronteira\nR$ 500,00",7.8,3.02,3.7,.65,22,MINT,True); text(s,"e o caminho excepcional:\nvalor negativo gera ValueError.",7.8,4.25,4.15,.55,16,INK)
text(s,"coverage erase  →  coverage run --branch -m pytest  →  coverage html",1.05,6.23,11.3,.28,14,BLUE,True,PP_ALIGN.CENTER,"Consolas")

# 9
s=base(9,"Relatório"); title(s,"Como ler o coverage report","Mostre o relatório real gerado no terminal e depois abra htmlcov/index.html.")
headers=["Stmts","Miss","Branch","BrPart","Cover"]
desc=["instruções\nconsideradas","linhas não\nexecutadas","caminhos de\ndecisão","caminhos\nparciais","percentual\ncalculado"]
for i,h in enumerate(headers):
 x=.65+i*2.5; box(s,x,2.35,2.12,2.45,PALE); text(s,h,x+.18,2.7,1.75,.3,17,BLUE,True,PP_ALIGN.CENTER); text(s,desc[i],x+.2,3.43,1.7,.55,15,INK,False,PP_ALIGN.CENTER)
text(s,"A maior evidência não é o número: é conseguir apontar, no HTML, qual linha ou caminho ainda falta testar.",.9,5.55,11.5,.5,17,NAVY,True,PP_ALIGN.CENTER)

# 10
s=base(10,"Limitação"); title(s,"100% de coverage ≠ ausência de bugs","Cobertura mede execução; asserções corretas medem comportamento.")
box(s,.75,2.0,5.8,3.65,RGBColor(255,242,242)); text(s,"REGRA COM BUG",1.15,2.42,2.6,.3,15,RED,True); text(s,"Cliente VIP recebe\n15% de desconto",1.15,3.02,4.55,.66,24,NAVY,True); text(s,"A regra esperada era 20%.",1.15,4.35,4.2,.3,16,MUTED)
box(s,6.95,2.0,5.65,3.65,PALE); text(s,"TESTE FRACO",7.35,2.42,2.6,.3,15,AMBER,True); text(s,"assert desconto_vip_com_bug(100) > 0",7.35,3.15,4.75,.35,15,NAVY,True,"Consolas"); text(s,"Ele passa e pode executar todas as linhas,\nmas nunca verifica o resultado correto: 80.",7.35,4.15,4.65,.62,16,INK)
text(s,"Cobertura alta + testes fracos = falsa sensação de segurança",1.1,6.1,11.1,.3,16,RED,True,PP_ALIGN.CENTER)

# 11
s=base(11,"Encerramento"); title(s,"A métrica é uma bússola, não o destino","O que levamos para o próximo projeto?")
for y,lead,body,c in [(2.05,"1","Use coverage para enxergar lacunas.",BLUE),(3.15,"2","Prefira branch coverage ao analisar decisões.",CYAN),(4.25,"3","Projete boas asserções e casos de fronteira.",MINT)]:
 box(s,1.0,y,.56,.56,c); text(s,lead,1.0,y+.12,.56,.18,13,NAVY,True,PP_ALIGN.CENTER); text(s,body,1.95,y+.08,8.7,.35,19,NAVY,True)
box(s,1.0,5.55,11.3,.7,NAVY); text(s,"Coverage mostra quanto foi exercitado. A qualidade dos testes mostra se aquilo foi verificado de forma significativa.",1.25,5.74,10.8,.25,16,WHITE,True,PP_ALIGN.CENTER)
text(s,"Referências: coverage.readthedocs.io · docs.pytest.org · pytest-cov.readthedocs.io",1.1,6.55,11,.25,10,MUTED,False,PP_ALIGN.CENTER)

prs.save(OUT)
print(f"Apresentação criada em: {OUT}")
